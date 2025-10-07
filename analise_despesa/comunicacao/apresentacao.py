# analise_despesa/comunicacao/apresentacao.py (VERSÃO FINAL COM GRÁFICOS E EXEMPLOS)

import logging
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from . import graficos_para_ppt # Importa o novo módulo de gráficos

logger = logging.getLogger(__name__)

def _add_slide_com_titulo(prs, titulo_texto):
    """Adiciona um novo slide com um título formatado."""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = titulo_texto
    title.text_frame.paragraphs[0].font.name = 'Calibri'
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    return slide

def _add_text_box_com_formatacao(slide, top, text_runs):
    """Adiciona uma caixa de texto com formatação customizada."""
    txBox = slide.shapes.add_textbox(Inches(1), top, Inches(7.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for text, size, bold, is_formula in text_runs:
        run = p.add_run()
        run.text = text
        run.font.name = 'Calibri'
        run.font.size = Pt(size)
        run.font.bold = bold
        if is_formula:
            run.font.color.rgb = RGBColor(0, 112, 192)

def gerar_pptx_metodologia(caminho_saida: str):
    """Cria uma apresentação PowerPoint (.pptx) explicando a metodologia do robô."""
    logger.info(f"Gerando apresentação da metodologia em: {caminho_saida}")
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # --- Diretório temporário para gráficos ---
    output_dir = os.path.dirname(caminho_saida)
    temp_chart_paths = []

    try:
        # --- Slide 1: Título ---
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title, subtitle = slide.shapes.title, slide.placeholders[1]
        title.text = "Metodologia e Análises do Robô de Despesas"
        subtitle.text = "Documentação Técnica e Conceitual Gerada Automaticamente"

        # --- Slide 2: Execução Orçamentária ---
        slide = _add_slide_com_titulo(prs, "Análise 1: Execução Orçamentária e Score de Criticidade")
        texto = [
            ("Pergunta de Negócio:", 16, True, False), (" Quais projetos representam maior risco e precisam de atenção imediata?\n\n", 16, False, False),
            ("Metodologia Aplicada:\n", 14, True, False), ("Combina a execução orçamentária (% gasto) com a detecção de ocorrências atípicas (ver Análise 4) para gerar um score de risco.\n\n", 14, False, False),
            ("Exemplo Prático:\n", 14, True, False),
            ("  - Projeto A: 95% de execução → ", 14, False, False), ("Score: ↑ Alto\n", 14, True, True),
            ("  - Projeto B: 80% de execução + Ocorrência Atípica → ", 14, False, False), ("Score: ↑ Alto\n", 14, True, True),
            ("  - Projeto C: 50% de execução + Ocorrência Atípica → ", 14, False, False), ("Score: → Médio", 14, True, True),
        ]
        _add_text_box_com_formatacao(slide, Inches(1.5), texto)
        caminho_grafico_score = os.path.join(output_dir, "temp_score.png")
        temp_chart_paths.append(graficos_para_ppt.criar_grafico_score(caminho_grafico_score))
        slide.shapes.add_picture(caminho_grafico_score, Inches(8.5), Inches(2.0), width=Inches(7))

        # --- Slide 3: Desempenho Mensal ---
        slide = _add_slide_com_titulo(prs, "Análise 2: Desempenho Mensal e Sinalização da IA")
        texto = [
            ("Pergunta de Negócio:", 16, True, False), (" O ritmo dos meus gastos mensais está normal ou houve algum mês com comportamento estranho?\n\n", 16, False, False),
            ("Metodologia Aplicada (IA - Isolation Forest):\n", 14, True, False),
            ("O modelo analisa a série de gastos totais da unidade e aprende o padrão de 'normalidade' do ano. Meses que fogem muito desse padrão são sinalizados, indicando desvios significativos no ritmo financeiro que merecem investigação.", 14, False, False),
        ]
        _add_text_box_com_formatacao(slide, Inches(1.5), texto)
        caminho_grafico_tendencia = os.path.join(output_dir, "temp_tendencia.png")
        temp_chart_paths.append(graficos_para_ppt.criar_grafico_tendencia(caminho_grafico_tendencia))
        slide.shapes.add_picture(caminho_grafico_tendencia, Inches(8.5), Inches(2.5), width=Inches(7))

        # --- Slide 4: Top 5 Fornecedores ---
        slide = _add_slide_com_titulo(prs, "Análise 3: Top 5 Fornecedores")
        texto = [
            ("Pergunta de Negócio:", 16, True, False), (" Para onde está indo o dinheiro? Quais são meus parceiros mais estratégicos?\n\n", 16, False, False),
            ("Metodologia Aplicada:\n", 14, True, False),
            ("É uma agregação direta que soma todos os valores pagos a cada fornecedor ao longo do ano. Os cinco maiores são exibidos. A análise é feita separadamente para projetos Exclusivos e Compartilhados, permitindo uma visão clara da dependência de fornecedores em cada contexto.", 14, False, False),
        ]
        _add_text_box_com_formatacao(slide, Inches(1.5), texto)

        # --- Slide 5: Ocorrências Atípicas ---
        slide = _add_slide_com_titulo(prs, "Análise 4: Ocorrências Atípicas de Contexto (IA Híbrida)")
        texto = [
            ("Pergunta de Negócio:", 16, True, False), (" Existem lançamentos que são estranhos do ponto de vista do negócio, mesmo com valor normal?\n\n", 16, False, False),
            ("Metodologia Aplicada (Z-Score + Isolation Forest):\n", 14, True, False),
            ("1. ", 14, True, False), ("Cálculo do Z-Score do Valor:", 14, False, False), (" Mede o quão atípico um valor é para sua própria combinação Fornecedor-Projeto.\n", 14, False, False),
            ("2. ", 14, True, False), ("Cálculo de Familiaridade:", 14, False, False), (" Mede a frequência do Fornecedor e do Projeto na unidade.\n", 14, False, False),
            ("3. ", 14, True, False), ("Treinamento da IA:", 14, False, False), (" O modelo é alimentado com [Z-Score, Freq. Fornecedor, Freq. Projeto] para encontrar combinações raras.", 14, False, False),
        ]
        _add_text_box_com_formatacao(slide, Inches(1.5), texto)
        caminho_grafico_contexto = os.path.join(output_dir, "temp_contexto.png")
        temp_chart_paths.append(graficos_para_ppt.criar_grafico_contexto(caminho_grafico_contexto))
        slide.shapes.add_picture(caminho_grafico_contexto, Inches(8.5), Inches(2.5), width=Inches(7))

        # --- Slide 6: Clusterização ---
        slide = _add_slide_com_titulo(prs, "Análise 5: Análise de Comportamento das Contas (IA)")
        texto = [
            ("Pergunta de Negócio:", 16, True, False), (" Qual é o 'DNA' dos meus tipos de despesa? São recorrentes, esporádicos, de alto ou baixo valor?\n\n", 16, False, False),
            ("Metodologia Aplicada (K-Means Clustering):\n", 14, True, False),
            ("1. ", 14, True, False), ("Extração do DNA:", 14, False, False), (" Para cada tipo de conta (Nível 4), calcula-se Valor Total, Frequência e o Coeficiente de Variação (CV), que mede a imprevisibilidade.\n", 14, False, False),
            ("   CV = (Desvio Padrão dos Gastos Mensais) / (Média dos Gastos Mensais)\n\n", 14, True, True),
            ("2. ", 14, True, False), ("Agrupamento (Clusterização):", 14, False, False), (" O modelo K-Means agrupa as contas com 'DNA' similar em perfis de comportamento.", 14, False, False),
        ]
        _add_text_box_com_formatacao(slide, Inches(1.5), texto)

        prs.save(caminho_saida)
        logger.info("Apresentação da metodologia gerada com sucesso.")

    finally:
        # --- Limpeza dos arquivos temporários ---
        for path in temp_chart_paths:
            if os.path.exists(path):
                os.remove(path)
                logger.debug(f"Arquivo de gráfico temporário removido: {path}")

